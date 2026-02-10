"""
支持子图分析的增强生成器
自动拆分子图并生成详细分析
"""

from pathlib import Path
from typing import Dict, List, Optional, Tuple
import sys

sys.path.insert(0, str(Path(__file__).parent))


class SubfigureAnalyzer:
    """子图分析器"""
    
    def __init__(self, llm, smart_agent):
        self.llm = llm
        self.smart_agent = smart_agent
    
    def ensure_subfigures_split(self, fig_num: int) -> Dict:
        """确保子图已拆分，如果没有则自动拆分"""
        
        if fig_num not in self.smart_agent.figure_map:
            raise ValueError(f"Figure {fig_num} 不存在")
        
        fig_data = self.smart_agent.figure_map[fig_num]
        
        # 检查是否已经拆分
        if "subfigures" in fig_data and fig_data["subfigures"]:
            print(f"[Analyzer] Figure {fig_num} 已拆分，有 {len(fig_data['subfigures'])} 个子图")
            return fig_data["subfigures"]
        
        # 如果没有拆分，尝试拆分
        if not fig_data.get("split_attempted"):
            print(f"[Analyzer] 正在拆分 Figure {fig_num}...")
            
            # 检查拆分器是否可用
            if not hasattr(self.smart_agent, 'splitter'):
                print(f"[Analyzer] ⚠️ 警告: splitter未初始化")
                fig_data["split_attempted"] = True
                return {}
            
            if not self.smart_agent.splitter.enabled:
                print(f"[Analyzer] ⚠️ 警告: figure-separator未安装或未启用")
                print(f"[Analyzer] 提示: 安装figure-separator以启用子图拆分功能")
                fig_data["split_attempted"] = True
                return {}
            
            try:
                output_dir = Path(self.smart_agent.data_dir) / "images"
                output_dir.mkdir(exist_ok=True)
                
                print(f"[Analyzer] 使用figure-separator拆分...")
                subfigs = self.smart_agent.splitter.split(
                    fig_data["path"],
                    str(output_dir),
                    fig_num,
                    use_numbers=True  # 使用数字标注 (1, 2, 3...)
                )
                
                if subfigs:
                    fig_data["subfigures"] = subfigs
                    fig_data["split_attempted"] = True
                    print(f"[Analyzer] ✓ Figure {fig_num} 拆分完成，共 {len(subfigs)} 个子图")
                    for label, path in subfigs.items():
                        print(f"[Analyzer]   - 子图 {fig_num}{label}: {path}")
                    return subfigs
                else:
                    print(f"[Analyzer] ⚠️ 警告: 未检测到子图")
                    fig_data["split_attempted"] = True
                    return {}
                
            except Exception as e:
                print(f"[Analyzer] ❌ Figure {fig_num} 拆分失败: {e}")
                import traceback
                traceback.print_exc()
                fig_data["split_attempted"] = True
                return {}
        
        return fig_data.get("subfigures", {})
    
    def analyze_subfigure(self, fig_num: int, subfig_label: str) -> Dict:
        """分析单个子图"""
        
        # 确保子图已拆分
        subfigures = self.ensure_subfigures_split(fig_num)
        
        if not subfigures or subfig_label not in subfigures:
            raise ValueError(f"子图 {fig_num}{subfig_label} 不存在")
        
        subfig_path = subfigures[subfig_label]
        
        print(f"[Analyzer] 分析子图 {fig_num}{subfig_label}...")
        
        # 1. 识别图表类型
        type_prompt = """请识别这张图的类型。从以下选项中选择最合适的：

热图(Heatmap)、散点图(Scatter)、柱状图(Bar)、折线图(Line)、
箱线图(Box)、小提琴图(Violin)、韦恩图(Venn)、
网络图(Network)、流程图(Flowchart)、架构图(Architecture)、
饼图(Pie)、雷达图(Radar)、树状图(Tree)、
密度图(Density)、直方图(Histogram)、其他(Other)

只需回答类型名称。"""
        
        chart_type = self.llm.analyze_image(subfig_path, type_prompt)
        
        # 2. 详细分析
        analysis_prompt = f"""这是Figure {fig_num}的子图{subfig_label}。

请详细分析这张图：
1. 图中展示了什么内容？
2. 横轴和纵轴分别代表什么？（如果有）
3. 主要的数据趋势或发现是什么？
4. 图中的颜色/标记代表什么含义？（如果有）
5. 这个子图想要说明什么结论？

请用清晰、易懂的语言回答，就像在给同事解释一样。"""
        
        analysis = self.llm.analyze_image(subfig_path, analysis_prompt)
        
        # 3. 获取文中对这个子图的描述
        context = self._get_subfigure_context(fig_num, subfig_label)
        
        return {
            "figure": fig_num,
            "subfigure": subfig_label,
            "path": subfig_path,
            "chart_type": chart_type.strip(),
            "analysis": analysis,
            "context": context
        }
    
    def analyze_all_subfigures(self, fig_num: int) -> List[Dict]:
        """分析一个Figure的所有子图"""
        
        # 确保子图已拆分
        subfigures = self.ensure_subfigures_split(fig_num)
        
        if not subfigures:
            return []
        
        results = []
        for subfig_label in sorted(subfigures.keys()):
            try:
                result = self.analyze_subfigure(fig_num, subfig_label)
                results.append(result)
            except Exception as e:
                print(f"[Analyzer] 分析子图 {fig_num}{subfig_label} 失败: {e}")
        
        return results
    
    def _get_subfigure_context(self, fig_num: int, subfig_label: str) -> str:
        """获取文中对子图的描述"""
        
        # 尝试匹配 "Figure 1a" 或 "Fig. 1a" 等格式
        patterns = [
            f"Figure {fig_num}{subfig_label}",
            f"Fig. {fig_num}{subfig_label}",
            f"figure {fig_num}{subfig_label}",
            f"fig. {fig_num}{subfig_label}",
        ]
        
        descriptions = []
        for text_data in self.smart_agent.texts:
            content = text_data["content"]
            for pattern in patterns:
                if pattern in content:
                    # 提取前后文
                    idx = content.find(pattern)
                    start = max(0, idx - 100)
                    end = min(len(content), idx + 200)
                    context = content[start:end]
                    descriptions.append(context)
                    break
            if descriptions:
                break
        
        return descriptions[0] if descriptions else ""


def generate_subfigure_report(
    analyzer: SubfigureAnalyzer,
    fig_num: int,
    output_format: str = "PDF"
) -> str:
    """生成子图分析报告"""
    
    from datetime import datetime
    
    print(f"[SubReport] 开始生成 Figure {fig_num} 子图报告...")
    
    # 分析所有子图
    subfig_analyses = analyzer.analyze_all_subfigures(fig_num)
    
    if not subfig_analyses:
        raise ValueError(f"Figure {fig_num} 没有子图或分析失败")
    
    # 准备报告内容
    content = {
        "概述": f"本报告详细分析了 Figure {fig_num} 的 {len(subfig_analyses)} 个子图。"
    }
    
    # 为每个子图创建章节
    for subfig in subfig_analyses:
        section_title = f"子图 {subfig['figure']}{subfig['subfigure']}"
        section_content = f"""**图表类型**: {subfig['chart_type']}

**详细分析**:
{subfig['analysis']}

**文中描述**:
{subfig['context'] if subfig['context'] else '（未找到明确的文中描述）'}
"""
        content[section_title] = section_content
    
    # 准备图片数据
    figures = []
    for subfig in subfig_analyses:
        figures.append({
            "number": f"{subfig['figure']}{subfig['subfigure']}",
            "path": subfig['path'],
            "description": f"{subfig['chart_type']}: {subfig['analysis'][:100]}..."
        })
    
    # 根据格式生成报告
    output_dir = Path(analyzer.smart_agent.data_dir)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    
    if output_format.upper() == "PDF":
        return _create_subfig_pdf(content, figures, output_dir, timestamp, fig_num)
    elif output_format.upper() in ["WORD", "DOCX"]:
        return _create_subfig_word(content, figures, output_dir, timestamp, fig_num)
    else:
        return _create_subfig_markdown(content, figures, output_dir, timestamp, fig_num)


def generate_subfigure_ppt(
    analyzer: SubfigureAnalyzer,
    fig_num: int
) -> str:
    """生成子图分析PPT"""
    
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from datetime import datetime
    
    print(f"[SubPPT] 开始生成 Figure {fig_num} 子图PPT...")
    
    # 分析所有子图
    subfig_analyses = analyzer.analyze_all_subfigures(fig_num)
    
    if not subfig_analyses:
        raise ValueError(f"Figure {fig_num} 没有子图或分析失败")
    
    # 创建PPT
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Figure {fig_num} 子图详细分析"
    subtitle.text = f"共 {len(subfig_analyses)} 个子图\n{datetime.now().strftime('%Y-%m-%d')}"
    
    # 为每个子图创建一页
    for subfig in subfig_analyses:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
        
        # 标题
        title = slide.shapes.title
        title.text = f"子图 {subfig['figure']}{subfig['subfigure']} - {subfig['chart_type']}"
        
        # 添加图片
        try:
            img_path = Path(subfig['path'])
            if img_path.exists():
                left = Inches(0.5)
                top = Inches(1.5)
                pic = slide.shapes.add_picture(
                    str(img_path),
                    left, top,
                    width=Inches(4.5)
                )
        except Exception as e:
            print(f"[SubPPT] 添加图片失败: {e}")
        
        # 添加分析文本框
        left = Inches(5.2)
        top = Inches(1.5)
        width = Inches(4.3)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # 添加分析内容
        analysis_text = subfig['analysis'][:400]  # 限制长度
        tf.text = analysis_text
        
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(12)
    
    # 保存
    output_path = Path(analyzer.smart_agent.data_dir) / f"subfigure_analysis_fig{fig_num}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(str(output_path))
    
    print(f"[SubPPT] ✅ PPT生成完成: {output_path}")
    return str(output_path)


def _create_subfig_pdf(content, figures, output_dir, timestamp, fig_num):
    """创建子图分析PDF"""
    
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib.enums import TA_CENTER
        
        output_path = output_dir / f"subfigure_analysis_fig{fig_num}_{timestamp}.pdf"
        
        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # 标题
        title_style = ParagraphStyle(
            'Title',
            parent=styles['Heading1'],
            fontSize=24,
            alignment=TA_CENTER,
            spaceAfter=30
        )
        story.append(Paragraph(f"Figure {fig_num} 子图详细分析", title_style))
        story.append(Spacer(1, 20))
        
        # 添加内容
        for section_title, section_text in content.items():
            story.append(Paragraph(section_title, styles['Heading2']))
            story.append(Spacer(1, 6))
            
            # 清理文本
            clean_text = section_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            clean_text = clean_text.replace('**', '')  # 移除markdown bold
            
            story.append(Paragraph(clean_text, styles['Normal']))
            story.append(Spacer(1, 12))
        
        # 添加图片
        story.append(PageBreak())
        story.append(Paragraph("子图展示", styles['Heading2']))
        story.append(Spacer(1, 12))
        
        for fig in figures:
            try:
                img_path = Path(fig['path'])
                if img_path.exists():
                    story.append(Paragraph(f"子图 {fig['number']}", styles['Heading3']))
                    img = Image(str(img_path), width=4*inch, height=None)
                    img.hAlign = 'CENTER'
                    story.append(img)
                    story.append(Spacer(1, 12))
            except Exception as e:
                print(f"添加图片失败: {e}")
        
        doc.build(story)
        return str(output_path)
        
    except ImportError:
        print("reportlab未安装，降级到Markdown")
        return _create_subfig_markdown(content, figures, output_dir, timestamp, fig_num)


def _create_subfig_word(content, figures, output_dir, timestamp, fig_num):
    """创建子图分析Word文档"""
    
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        output_path = output_dir / f"subfigure_analysis_fig{fig_num}_{timestamp}.docx"
        
        doc = Document()
        
        # 标题
        title = doc.add_heading(f"Figure {fig_num} 子图详细分析", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()
        
        # 添加内容
        for section_title, section_text in content.items():
            doc.add_heading(section_title, level=1)
            doc.add_paragraph(section_text)
            doc.add_paragraph()
        
        # 添加图片
        doc.add_heading("子图展示", level=1)
        
        for fig in figures:
            try:
                img_path = Path(fig['path'])
                if img_path.exists():
                    doc.add_heading(f"子图 {fig['number']}", level=2)
                    doc.add_picture(str(img_path), width=Inches(4))
                    doc.add_paragraph()
            except Exception as e:
                print(f"添加图片失败: {e}")
        
        doc.save(str(output_path))
        return str(output_path)
        
    except ImportError:
        print("python-docx未安装，降级到Markdown")
        return _create_subfig_markdown(content, figures, output_dir, timestamp, fig_num)


def _create_subfig_markdown(content, figures, output_dir, timestamp, fig_num):
    """创建子图分析Markdown包"""
    
    import shutil
    
    # 创建文件夹
    report_folder = output_dir / f"subfigure_analysis_fig{fig_num}_{timestamp}"
    report_folder.mkdir(exist_ok=True)
    
    images_folder = report_folder / "images"
    images_folder.mkdir(exist_ok=True)
    
    # 复制图片
    copied_figures = []
    for fig in figures:
        try:
            src_path = Path(fig['path'])
            if src_path.exists():
                dst_filename = f"subfig_{fig['number']}{src_path.suffix}"
                dst_path = images_folder / dst_filename
                shutil.copy2(src_path, dst_path)
                
                copied_figures.append({
                    "number": fig["number"],
                    "path": f"images/{dst_filename}",
                    "description": fig["description"]
                })
        except Exception as e:
            print(f"复制图片失败: {e}")
    
    # 生成Markdown
    md_content = f"# Figure {fig_num} 子图详细分析\n\n"
    
    for section_title, section_text in content.items():
        md_content += f"## {section_title}\n\n{section_text}\n\n"
    
    md_content += "## 子图展示\n\n"
    for fig in copied_figures:
        md_content += f"### 子图 {fig['number']}\n\n"
        md_content += f"![{fig['number']}]({fig['path']})\n\n"
        md_content += f"{fig['description']}\n\n"
    
    # 保存
    md_path = report_folder / "analysis.md"
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(md_content)
    
    return str(report_folder)


__all__ = ["SubfigureAnalyzer", "generate_subfigure_report", "generate_subfigure_ppt"]
