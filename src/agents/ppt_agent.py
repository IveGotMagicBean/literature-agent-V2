"""
增强版PPT生成Agent
修复图片插入问题，支持更多自定义选项
"""

from pathlib import Path
from typing import Dict, List, Optional
import sys

sys.path.insert(0, str(Path(__file__).parent))


class EnhancedPPTAgent:
    """增强版PPT生成Agent"""
    
    def __init__(self, llm, smart_agent):
        self.llm = llm
        self.smart_agent = smart_agent
    
    def generate(
        self, 
        template: str = "学术风格", 
        language: str = "中文",
        include_figures: bool = True,
        max_figures: int = 5
    ) -> Optional[str]:
        """
        生成PPT
        
        Args:
            template: PPT模板（学术风格/简约风格/商务风格）
            language: 语言（中文/英文/双语）
            include_figures: 是否包含图片
            max_figures: 最多包含多少张图片
        
        Returns:
            PPT文件路径
        """
        
        # 检查是否已加载文献
        if not self.smart_agent.texts:
            raise ValueError("请先上传PDF文献")
        
        print("[PPT Agent] 开始生成PPT...")
        
        # 1. 提取论文结构
        print("[PPT Agent] 提取论文结构...")
        structure = self._extract_structure()
        
        # 2. 提取关键图表
        figures = []
        if include_figures and self.smart_agent.figure_map:
            print("[PPT Agent] 提取关键图表...")
            figures = self._extract_key_figures(max_figures)
        
        # 3. 生成每一页内容
        print("[PPT Agent] 生成幻灯片内容...")
        slides = self._generate_slides(structure, figures, language)
        
        # 4. 创建PPT文件
        print("[PPT Agent] 创建PPT文件...")
        ppt_path = self._create_pptx(slides, template, language)
        
        print(f"[PPT Agent] ✅ PPT生成完成: {ppt_path}")
        return ppt_path
    
    def _extract_structure(self) -> Dict:
        """提取论文结构"""
        
        prompt = """请分析这篇论文的结构，提取以下信息：

1. 标题
2. 作者和单位
3. 研究背景（1-2句话）
4. 研究问题/动机
5. 主要方法
6. 关键实验
7. 主要结果
8. 结论

文献内容：
{text}

请用JSON格式返回，格式如下：
{{
    "title": "...",
    "authors": "...",
    "background": "...",
    "motivation": "...",
    "method": "...",
    "experiments": "...",
    "results": "...",
    "conclusion": "..."
}}
"""
        
        # 取前5000字
        text_sample = self.smart_agent.full_text[:5000]
        
        response = self.llm.chat(prompt.format(text=text_sample))
        
        # 解析JSON
        import json
        try:
            structure = json.loads(response)
        except:
            # 如果解析失败，返回默认结构
            structure = {
                "title": "学术论文",
                "authors": "作者",
                "background": "研究背景",
                "motivation": "研究动机",
                "method": "方法",
                "experiments": "实验",
                "results": "结果",
                "conclusion": "结论"
            }
        
        return structure
    
    def _extract_key_figures(self, max_figures: int = 5) -> List[Dict]:
        """提取关键图表"""
        
        key_figures = []
        
        # 取前N张重要图表
        for fig_num, fig_data in list(self.smart_agent.figure_map.items())[:max_figures]:
            if not fig_data or "path" not in fig_data:
                continue
            
            # 确保图片路径存在
            fig_path = Path(fig_data["path"])
            if not fig_path.exists():
                print(f"[PPT Agent] 警告: 图片不存在 {fig_path}")
                continue
            
            # 获取图表描述
            description = self._get_figure_summary(fig_num)
            
            key_figures.append({
                "number": fig_num,
                "path": str(fig_path.absolute()),  # 使用绝对路径
                "description": description
            })
        
        return key_figures
    
    def _get_figure_summary(self, fig_num: int) -> str:
        """获取图表简介"""
        
        # 从文中提取描述
        descriptions = []
        for text_data in self.smart_agent.texts:
            refs = self.smart_agent.ref_extractor.extract_references(text_data["content"])
            for ref in refs:
                if ref["figure"] == fig_num:
                    context = self.smart_agent.ref_extractor.get_context(
                        text_data["content"], ref, 150
                    )
                    descriptions.append(context)
                    if len(descriptions) >= 1:
                        break
            if descriptions:
                break
        
        return descriptions[0] if descriptions else f"Figure {fig_num}"
    
    def _generate_slides(self, structure: Dict, figures: List[Dict], language: str) -> List[Dict]:
        """生成每一页的内容"""
        
        slides = []
        
        # 第1页：标题页
        slides.append({
            "type": "title",
            "title": structure.get("title", "学术论文"),
            "subtitle": structure.get("authors", "作者")
        })
        
        # 第2页：研究背景
        slides.append({
            "type": "content",
            "title": "研究背景" if language == "中文" else "Background",
            "content": structure.get("background", "研究背景")
        })
        
        # 第3页：研究动机
        slides.append({
            "type": "content",
            "title": "研究动机" if language == "中文" else "Motivation",
            "content": structure.get("motivation", "研究动机")
        })
        
        # 第4页：方法
        slides.append({
            "type": "content",
            "title": "方法" if language == "中文" else "Method",
            "content": structure.get("method", "方法")
        })
        
        # 第5-N页：关键图表
        for fig in figures:
            slides.append({
                "type": "figure",
                "title": f"Figure {fig['number']}",
                "image": fig["path"],
                "description": fig["description"]
            })
        
        # 倒数第2页：结果
        slides.append({
            "type": "content",
            "title": "主要结果" if language == "中文" else "Results",
            "content": structure.get("results", "结果")
        })
        
        # 最后一页：结论
        slides.append({
            "type": "content",
            "title": "结论" if language == "中文" else "Conclusion",
            "content": structure.get("conclusion", "结论")
        })
        
        return slides
    
    def _create_pptx(self, slides: List[Dict], template: str, language: str) -> str:
        """创建PPT文件 - 使用python-pptx直接生成"""
        
        from datetime import datetime
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        print("[PPT Agent] 创建PPT...")
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 主题颜色
        theme_colors = {
            "学术风格": {"primary": RGBColor(31, 78, 121), "secondary": RGBColor(68, 114, 196)},
            "商务风格": {"primary": RGBColor(68, 84, 106), "secondary": RGBColor(112, 128, 144)},
            "科技风格": {"primary": RGBColor(46, 125, 50), "secondary": RGBColor(102, 187, 106)}
        }
        colors = theme_colors.get(template, theme_colors["学术风格"])
        
        for idx, slide_data in enumerate(slides):
            print(f"[PPT Agent] 生成第 {idx+1}/{len(slides)} 页...")
            
            if slide_data["type"] == "title":
                # 标题页
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
                
                # 标题
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
                title_frame = title_box.text_frame
                title = title_frame.add_paragraph()
                title.text = slide_data["title"]
                title.font.size = Pt(44)
                title.font.bold = True
                title.font.color.rgb = colors["primary"]
                title.alignment = PP_ALIGN.CENTER
                
                # 副标题
                if slide_data.get("subtitle"):
                    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
                    subtitle_frame = subtitle_box.text_frame
                    subtitle = subtitle_frame.add_paragraph()
                    subtitle.text = slide_data["subtitle"]
                    subtitle.font.size = Pt(20)
                    subtitle.font.color.rgb = colors["secondary"]
                    subtitle.alignment = PP_ALIGN.CENTER
            
            elif slide_data["type"] == "content":
                # 内容页
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 标题
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
                title_frame = title_box.text_frame
                title = title_frame.add_paragraph()
                title.text = slide_data["title"]
                title.font.size = Pt(32)
                title.font.bold = True
                title.font.color.rgb = colors["primary"]
                
                # 内容
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                # 解析内容（支持列表）
                content_text = slide_data["content"]
                for line in content_text.split('\n'):
                    line = line.strip()
                    if not line:
                        continue
                    
                    p = content_frame.add_paragraph()
                    
                    # 检测是否是列表项
                    if line.startswith('- ') or line.startswith('• '):
                        p.text = line[2:]
                        p.level = 0
                        p.font.size = Pt(18)
                    elif line.startswith('  - ') or line.startswith('  • '):
                        p.text = line[4:]
                        p.level = 1
                        p.font.size = Pt(16)
                    else:
                        p.text = line
                        p.font.size = Pt(18)
                    
                    p.font.color.rgb = RGBColor(64, 64, 64)
                    p.space_after = Pt(6)
            
            elif slide_data["type"] == "image":
                # 图片页
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 标题
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
                title_frame = title_box.text_frame
                title = title_frame.add_paragraph()
                title.text = slide_data["title"]
                title.font.size = Pt(32)
                title.font.bold = True
                title.font.color.rgb = colors["primary"]
                
                # 图片
                image_path = slide_data.get("image_path")
                if image_path and Path(image_path).exists():
                    try:
                        # 添加图片，居中显示
                        left = Inches(1.5)
                        top = Inches(1.8)
                        height = Inches(4.5)
                        slide.shapes.add_picture(str(image_path), left, top, height=height)
                    except Exception as e:
                        print(f"[PPT Agent] 添加图片失败: {e}")
                
                # 说明文字
                if slide_data.get("description"):
                    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.8))
                    desc_frame = desc_box.text_frame
                    desc = desc_frame.add_paragraph()
                    desc.text = slide_data["description"]
                    desc.font.size = Pt(14)
                    desc.font.color.rgb = RGBColor(96, 96, 96)
                    desc.alignment = PP_ALIGN.CENTER
        
        # 保存
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = Path("data") / f"presentation_{timestamp}.pptx"
        output_path.parent.mkdir(exist_ok=True)
        
        prs.save(str(output_path))
        print(f"[PPT Agent] ✅ PPT已保存: {output_path}")
        
        return str(output_path)


__all__ = ["EnhancedPPTAgent"]
