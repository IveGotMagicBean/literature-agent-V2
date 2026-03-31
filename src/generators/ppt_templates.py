"""
美化版PPT生成器 - 专业配色和设计
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from datetime import datetime


class BeautifulPPTGenerator:
    """美化版PPT生成器"""
    
    # 配色方案
    THEMES = {
        "学术蓝": {
            "primary": RGBColor(31, 78, 121),      # 深蓝色
            "secondary": RGBColor(68, 114, 196),   # 中蓝色
            "accent": RGBColor(91, 155, 213),      # 浅蓝色
            "background": RGBColor(255, 255, 255), # 白色
            "text": RGBColor(0, 0, 0),             # 黑色
            "subtitle": RGBColor(89, 89, 89),      # 灰色
        },
        "科技绿": {
            "primary": RGBColor(0, 128, 96),       # 深绿色
            "secondary": RGBColor(0, 176, 80),     # 绿色
            "accent": RGBColor(146, 208, 80),      # 浅绿色
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(0, 0, 0),
            "subtitle": RGBColor(89, 89, 89),
        },
        "商务灰": {
            "primary": RGBColor(68, 68, 68),       # 深灰
            "secondary": RGBColor(127, 127, 127),  # 中灰
            "accent": RGBColor(191, 191, 191),     # 浅灰
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(0, 0, 0),
            "subtitle": RGBColor(89, 89, 89),
        }
    }
    
    def __init__(self, theme="学术蓝"):
        self.theme = self.THEMES.get(theme, self.THEMES["学术蓝"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title: str, subtitle: str):
        """添加标题页 - 美化版"""
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        # 背景色块（顶部）
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=0, top=0,
            width=self.prs.slide_width,
            height=Inches(2.5)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.theme["primary"]
        top_bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.8),
            width=Inches(9),
            height=Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(3.5),
            width=Inches(8),
            height=Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = self.theme["subtitle"]
        
        # 日期
        date_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(6.5),
            width=Inches(8),
            height=Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%Y年%m月%d日")
        
        p = date_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = self.theme["subtitle"]
        
        # 装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(3),
            top=Inches(4.7),
            width=Inches(4),
            height=Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["accent"]
        line.line.fill.background()
    
    def add_content_slide(self, title: str, content: str, page_num: int = None):
        """添加内容页 - 美化版"""
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 顶部色块
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=0, top=0,
            width=self.prs.slide_width,
            height=Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["primary"]
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.15),
            width=Inches(8),
            height=Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 页码
        if page_num:
            page_box = slide.shapes.add_textbox(
                left=Inches(9),
                top=Inches(0.2),
                width=Inches(0.5),
                height=Inches(0.4)
            )
            page_frame = page_box.text_frame
            page_frame.text = str(page_num)
            p = page_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 内容框
        content_box = slide.shapes.add_textbox(
            left=Inches(0.8),
            top=Inches(1.5),
            width=Inches(8.4),
            height=Inches(5.3)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = content
        
        # 格式化内容
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.theme["text"]
            paragraph.space_after = Pt(12)
        
        # 左侧装饰条
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0.3),
            top=Inches(1.2),
            width=Inches(0.15),
            height=Inches(5.8)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.theme["accent"]
        accent_bar.line.fill.background()
    
    def add_figure_slide(self, title: str, image_path: str, description: str, page_num: int = None):
        """添加图片页 - 美化版"""
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 顶部色块
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=0, top=0,
            width=self.prs.slide_width,
            height=Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["primary"]
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.15),
            width=Inches(8),
            height=Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 页码
        if page_num:
            page_box = slide.shapes.add_textbox(
                left=Inches(9),
                top=Inches(0.2),
                width=Inches(0.5),
                height=Inches(0.4)
            )
            page_frame = page_box.text_frame
            page_frame.text = str(page_num)
            p = page_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 图片
        try:
            img_path = Path(image_path)
            if img_path.exists():
                # 图片框（带阴影效果）
                left = Inches(1.5)
                top = Inches(1.3)
                width = Inches(7)
                
                pic = slide.shapes.add_picture(
                    str(img_path),
                    left, top,
                    width=width
                )
                
                # 添加阴影
                pic.shadow.inherit = False
        except Exception as e:
            print(f"[PPT] 添加图片失败: {e}")
        
        # 描述文本框
        desc_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(6.5),
            width=Inches(8),
            height=Inches(0.8)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = description[:150] + "..." if len(description) > 150 else description
        desc_frame.word_wrap = True
        
        p = desc_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = self.theme["subtitle"]
    
    def save(self, output_path: str):
        """保存PPT"""
        self.prs.save(output_path)
        return output_path


__all__ = ["BeautifulPPTGenerator"]
